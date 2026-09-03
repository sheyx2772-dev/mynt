#!/usr/bin/env python3
"""Builds the Flex pitch deck, in English and in Uzbek.

One layout, two dictionaries. The alternative — writing the deck twice — is two
decks that drift apart, and the Uzbek one always drifts second.

Everything factual here is checked against the product: the prices come from the
same constants the checkout uses, the screenshots are of the running site, and
the status slide says what actually works rather than what is planned. A deck
that overstates is found out in the meeting after it.

Run scripts/deck_frames.py first. Writes to deck/.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FRAMES = ROOT / "deck" / "frames"
OUT = ROOT / "deck"

# The brand, and nothing beside it. One accent used sparingly is the whole
# reason the lime reads as emphasis rather than decoration.
INK = RGBColor(0x0E, 0x0A, 0x1B)
LIME = RGBColor(0xAB, 0xFF, 0x09)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xFA, 0xFA, 0xF8)
# Greys biased toward the ink's violet rather than pure neutral, so nothing on a
# slide looks like it came from a different palette.
DIM = RGBColor(0x8B, 0x87, 0x94)
DIM_DARK = RGBColor(0x6B, 0x66, 0x78)
LIME_INK = RGBColor(0x5C, 0x7A, 0x0A)

# Arial, because a deck is opened on somebody else's laptop and a font that is
# not there is silently replaced by one that ruins the layout. The design is
# carried by size, space and colour instead.
FONT = "Arial"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.85)


def text(slide, s, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0

    for i, line in enumerate(str(s).split("\n")):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = FONT
        run.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def slide(prs, dark=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, INK if dark else PAPER)
    return s


def eyebrow(s, label, dark=True):
    """A lime rule and a small word: the same header on every content slide."""
    rect(s, M, Inches(0.72), Inches(0.42), Pt(4), LIME)
    text(s, label.upper(), M + Inches(0.62), Inches(0.62), Inches(6), Inches(0.3),
         12, DIM if dark else DIM_DARK, bold=True)


def title(s, t, dark=True, y=1.15, size=40):
    text(s, t, M, Inches(y), W - M * 2, Inches(1.6), size, WHITE if dark else INK,
         bold=True, spacing=0.95)


def footer(s, page, dark=True):
    text(s, "flex.com.uz", M, H - Inches(0.62), Inches(3), Inches(0.3), 10,
         DIM if dark else DIM_DARK)
    text(s, str(page), W - M - Inches(1), H - Inches(0.62), Inches(1), Inches(0.3), 10,
         DIM if dark else DIM_DARK, align=PP_ALIGN.RIGHT)


def picture(s, name, x, y, height):
    path = FRAMES / f"{name}.png"
    return s.shapes.add_picture(str(path), x, y, height=height)


# ---------------------------------------------------------------- the slides


def build(lang, t):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    page = 0

    # 1 — title
    s = slide(prs)
    rect(s, 0, 0, Inches(0.18), H, LIME)
    text(s, "flex", M, Inches(2.15), Inches(6), Inches(1.1), 60, WHITE, bold=True)
    text(s, t["tagline"], M, Inches(3.3), Inches(8.4), Inches(1.6), 30, LIME,
         bold=True, spacing=1.0)
    text(s, t["sub"], M, Inches(4.9), Inches(8.4), Inches(1.2), 15, DIM, spacing=1.35)
    text(s, t["event"], M, H - Inches(1.15), Inches(8), Inches(0.4), 12, DIM, bold=True)

    # 2 — problem
    page += 1
    s = slide(prs, dark=False)
    eyebrow(s, t["nav_problem"], dark=False)
    title(s, t["problem_title"], dark=False)
    col = (W - M * 2 - Inches(0.9)) / 3
    for i, (head, body) in enumerate(t["problems"]):
        x = M + (col + Inches(0.45)) * i
        rect(s, x, Inches(3.0), Inches(0.34), Pt(4), LIME)
        text(s, head, x, Inches(3.35), col, Inches(0.8), 20, INK, bold=True, spacing=1.1)
        text(s, body, x, Inches(4.35), col, Inches(2), 14, DIM_DARK, spacing=1.45)
    footer(s, page, dark=False)

    # 3 — what it is
    page += 1
    s = slide(prs)
    eyebrow(s, t["nav_what"])
    title(s, t["what_title"])
    text(s, t["what_body"], M, Inches(3.0), Inches(6.4), Inches(2.4), 16, DIM, spacing=1.5)
    for i, step in enumerate(t["steps"]):
        y = Inches(4.7 + i * 0.62)
        text(s, f"{i + 1}", M, y, Inches(0.4), Inches(0.4), 15, LIME, bold=True)
        text(s, step, M + Inches(0.45), y, Inches(5.9), Inches(0.5), 15, WHITE)
    picture(s, "menu", W - M - Inches(2.75), Inches(0.95), Inches(6.1))
    footer(s, page)

    # 4 — two products
    page += 1
    s = slide(prs, dark=False)
    eyebrow(s, t["nav_products"], dark=False)
    title(s, t["products_title"], dark=False)
    half = (W - M * 2 - Inches(0.5)) / 2
    for i, (head, sub, price, bullets, is_lime) in enumerate(t["products"]):
        x = M + (half + Inches(0.5)) * i
        card = rect(s, x, Inches(2.75), half, Inches(3.75), LIME if is_lime else INK,
                    MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.055
        fg = INK if is_lime else WHITE
        mid = LIME_INK if is_lime else DIM
        text(s, head, x + Inches(0.5), Inches(3.15), half - Inches(1), Inches(0.5), 24, fg, bold=True)
        text(s, sub, x + Inches(0.5), Inches(3.72), half - Inches(1), Inches(0.4), 13, mid)
        text(s, price, x + Inches(0.5), Inches(4.25), half - Inches(1), Inches(0.6), 26,
             INK if is_lime else LIME, bold=True)
        text(s, "\n".join(f"·  {b}" for b in bullets), x + Inches(0.5), Inches(5.05),
             half - Inches(1), Inches(1.3), 13, fg if is_lime else DIM, spacing=1.5)
    footer(s, page, dark=False)

    # 5 — the venue loop
    page += 1
    s = slide(prs)
    eyebrow(s, t["nav_loop"])
    title(s, t["loop_title"])
    picture(s, "menu", M, Inches(2.5), Inches(4.15))
    picture(s, "counter", Inches(3.95), Inches(2.5), Inches(4.15))
    lx = Inches(7.55)
    for i, step in enumerate(t["loop"]):
        y = Inches(2.6 + i * 0.86)
        rect(s, lx, y + Inches(0.09), Inches(0.2), Pt(3), LIME)
        text(s, step, lx + Inches(0.4), y, W - lx - M - Inches(0.4), Inches(0.75), 14,
             WHITE, spacing=1.35)
    footer(s, page)

    # 6 — what a venue receives
    page += 1
    s = slide(prs, dark=False)
    eyebrow(s, t["nav_gets"], dark=False)
    title(s, t["gets_title"], dark=False)
    s.shapes.add_picture(str(FRAMES / "cards.png"), M, Inches(2.55), width=Inches(5.5))
    picture(s, "report", W - M - Inches(2.85), Inches(2.35), Inches(4.05))
    text(s, t["gets_body"], M, Inches(6.15), Inches(5.5), Inches(1.0), 12, DIM_DARK, spacing=1.4)
    footer(s, page, dark=False)

    # 7 — business model
    page += 1
    s = slide(prs)
    eyebrow(s, t["nav_model"])
    title(s, t["model_title"])
    for i, (band, price, note) in enumerate(t["bands"]):
        y = Inches(2.85 + i * 1.02)
        rect(s, M, y + Inches(0.42), Inches(0.24), Pt(3), LIME)
        text(s, band, M + Inches(0.5), y, Inches(4.6), Inches(0.5), 17, WHITE, bold=True)
        text(s, note, M + Inches(0.5), y + Inches(0.42), Inches(4.6), Inches(0.4), 12, DIM)
        text(s, price, M + Inches(4.6), y, Inches(3.05), Inches(0.6), 18, LIME, bold=True,
             align=PP_ALIGN.RIGHT)
    box = rect(s, W - M - Inches(3.6), Inches(2.85), Inches(3.6), Inches(3.35),
               RGBColor(0x1A, 0x16, 0x2A), MSO_SHAPE.ROUNDED_RECTANGLE)
    box.adjustments[0] = 0.07
    text(s, t["model_side_head"], W - M - Inches(3.15), Inches(3.25), Inches(2.7),
         Inches(0.4), 14, LIME, bold=True)
    text(s, t["model_side"], W - M - Inches(3.15), Inches(3.85), Inches(2.7), Inches(2.1),
         12, DIM, spacing=1.5)
    footer(s, page)

    # 8 — the arithmetic
    page += 1
    s = slide(prs, dark=False)
    eyebrow(s, t["nav_math"], dark=False)
    title(s, t["math_title"], dark=False)
    third = (W - M * 2 - Inches(0.9)) / 3
    for i, (venues, monthly, yearly) in enumerate(t["scenarios"]):
        x = M + (third + Inches(0.45)) * i
        card = rect(s, x, Inches(2.9), third, Inches(2.5), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.06
        card.line.color.rgb = RGBColor(0xE4, 0xE2, 0xE8)
        card.line.width = Pt(1)
        text(s, venues, x + Inches(0.42), Inches(3.25), third - Inches(0.84), Inches(0.5),
             15, DIM_DARK, bold=True)
        text(s, monthly, x + Inches(0.42), Inches(3.85), third - Inches(0.84), Inches(0.7),
             30, INK, bold=True)
        text(s, yearly, x + Inches(0.42), Inches(4.65), third - Inches(0.84), Inches(0.5),
             13, DIM_DARK)
    text(s, t["math_note"], M, Inches(5.85), W - M * 2, Inches(0.9), 12, DIM_DARK, spacing=1.45)
    footer(s, page, dark=False)

    # 9 — why now
    page += 1
    s = slide(prs)
    eyebrow(s, t["nav_why"])
    title(s, t["why_title"])
    for i, (head, body) in enumerate(t["why"]):
        y = Inches(2.9 + i * 1.32)
        text(s, head, M, y, Inches(4.4), Inches(0.5), 18, LIME, bold=True)
        text(s, body, M + Inches(4.8), y, W - M * 2 - Inches(4.8), Inches(1.1), 14, DIM,
             spacing=1.45)
    footer(s, page)

    # 10 — where it stands
    page += 1
    s = slide(prs, dark=False)
    eyebrow(s, t["nav_status"], dark=False)
    title(s, t["status_title"], dark=False)
    left = t["status_live"]
    for i, item in enumerate(left):
        y = Inches(2.7 + i * 0.47)
        # A drawn rule rather than a tick: U+2713 is not in Arial on Windows,
        # and a deck that shows a box instead of a tick has lost the room.
        rect(s, M, y + Inches(0.11), Inches(0.16), Pt(3), RGBColor(0x4D, 0x7C, 0x0F))
        text(s, item, M + Inches(0.38), y, Inches(6.0), Inches(0.45), 13, INK, spacing=1.25)
    box = rect(s, W - M - Inches(4.5), Inches(2.6), Inches(4.5), Inches(3.3), INK,
               MSO_SHAPE.ROUNDED_RECTANGLE)
    box.adjustments[0] = 0.055
    text(s, t["status_next_head"], W - M - Inches(4.05), Inches(3.0), Inches(3.6),
         Inches(0.4), 15, LIME, bold=True)
    text(s, "\n".join(f"·  {x}" for x in t["status_next"]), W - M - Inches(4.05),
         Inches(3.6), Inches(3.6), Inches(2.1), 12, DIM, spacing=1.55)
    footer(s, page, dark=False)

    # 11 — competition
    page += 1
    s = slide(prs)
    eyebrow(s, t["nav_comp"])
    title(s, t["comp_title"])
    for i, (who, what, gap) in enumerate(t["competitors"]):
        y = Inches(2.9 + i * 1.15)
        text(s, who, M, y, Inches(3.0), Inches(0.5), 17, WHITE, bold=True)
        text(s, what, M + Inches(3.2), y, Inches(4.0), Inches(0.9), 13, DIM, spacing=1.35)
        text(s, gap, M + Inches(7.5), y, W - M * 2 - Inches(7.5), Inches(0.9), 13, LIME,
             spacing=1.35)
    footer(s, page)

    # 12 — team
    page += 1
    s = slide(prs, dark=False)
    eyebrow(s, t["nav_team"], dark=False)
    title(s, t["team_title"], dark=False)
    third = (W - M * 2 - Inches(0.9)) / 3
    for i, (name, role, place) in enumerate(t["team"]):
        x = M + (third + Inches(0.45)) * i
        rect(s, x, Inches(3.0), Inches(0.3), Pt(4), LIME)
        text(s, name, x, Inches(3.35), third, Inches(0.8), 19, INK, bold=True, spacing=1.15)
        text(s, role, x, Inches(4.3), third, Inches(0.4), 14, LIME_INK, bold=True)
        text(s, place, x, Inches(4.75), third, Inches(0.4), 13, DIM_DARK)
    footer(s, page, dark=False)

    # 13 — ask
    page += 1
    s = slide(prs)
    rect(s, 0, 0, Inches(0.18), H, LIME)
    eyebrow(s, t["nav_ask"])
    title(s, t["ask_title"])
    text(s, t["ask_amount"], M, Inches(2.55), Inches(7.0), Inches(1.0), 38, LIME, bold=True)
    text(s, t["ask_stage"], M, Inches(3.5), Inches(6.4), Inches(0.5), 14, DIM)
    for i, (share, what) in enumerate(t["use"]):
        y = Inches(4.35 + i * 0.6)
        text(s, share, M, y, Inches(1.0), Inches(0.4), 15, WHITE, bold=True)
        text(s, what, M + Inches(1.2), y, Inches(5.2), Inches(0.5), 14, DIM)
    box = rect(s, W - M - Inches(4.0), Inches(2.5), Inches(4.0), Inches(3.3),
               RGBColor(0x1A, 0x16, 0x2A), MSO_SHAPE.ROUNDED_RECTANGLE)
    box.adjustments[0] = 0.055
    text(s, t["contact_head"], W - M - Inches(3.55), Inches(2.9), Inches(3.1), Inches(0.4),
         14, LIME, bold=True)
    text(s, t["contact"], W - M - Inches(3.55), Inches(3.5), Inches(3.1), Inches(2.1), 13,
         WHITE, spacing=1.7)
    footer(s, page)

    path = OUT / f"flex-pitch-{lang}.pptx"
    prs.save(path)
    size = path.stat().st_size / 1024
    print(f"  {path.name:<26} {len(prs.slides.__iter__.__self__._sldIdLst)} slides, {size:.0f} KB")


EN = {
    "tagline": "Your number is your identity.",
    "sub": "A unique number and an NFC object that opens it.\nOne tap. No app to install.",
    "event": "ICT WEEK UZBEKISTAN  ·  2026  ·  PRE-SEED",
    "nav_problem": "The problem",
    "problem_title": "Paper is printed once.\nEverything else changes.",
    "problems": [
        ("A card is out of date\nthe day it is printed",
         "A new job, a new number, a new company — and the box of cards in the drawer is wrong. People hand them over anyway."),
        ("A menu is a photograph\nof last month's prices",
         "A cafe reprints, or crosses lines out with a pen. The dish that ran out at lunch is still on every table."),
        ("A table cannot\nask for anything",
         "A guest raises a hand and waits for someone to look. The cafe never learns which tables waited or for how long."),
    ],
    "nav_what": "What Flex is",
    "what_title": "A number, an object,\nand a page behind it.",
    "what_body": "Three letters and three digits — MYN042 — become an address. A card, a ring, a bracelet or a sticker on a table opens it.",
    "steps": [
        "The guest taps their own phone against it",
        "The page opens in their own language",
        "Nobody installs anything, ever",
    ],
    "nav_products": "Two products",
    "products_title": "One platform, two customers.",
    "products": [
        ("Personal", "A number that is yours for life", "from 100,000 UZS",
         ["Profile, links, contacts", "Card, ring or bracelet", "Priced by rarity, paid once"], True),
        ("Venue", "Cafes, hotels and shops", "from 149,000 UZS / month",
         ["Menu and service list", "A tag for every table", "Calls, counter screen, report"], False),
    ],
    "nav_loop": "How a venue works",
    "loop_title": "The half a printed menu cannot do.",
    "loop": [
        "Every table has its own QR and its own NFC tag",
        "The guest taps it — the menu opens in their language",
        "One tap calls a waiter, asks for the bill, or leaves a review",
        "The request lands on the phone by the till, with the table number, and it rings",
        "Nobody signs in: the till phone opens one link and stays on it",
    ],
    "nav_gets": "What a venue receives",
    "gets_title": "Printed on Monday, answering by Tuesday.",
    "gets_body": "The owner types the tables, prints the sheet and puts a card on each one. Every code carries its own table number — so a request always says where it came from.",
    "nav_model": "Business model",
    "model_title": "A number is bought once.\nA venue pays every month.",
    "bands": [
        ("Personal number", "from 100,000 UZS", "Priced by how rare the combination is · paid once"),
        ("Device", "200,000 – 350,000 UZS", "Card, ring or bracelet · paid once"),
        ("Venue, up to 15 points", "149,000 UZS / mo", "A cafe with twelve tables"),
        ("Venue, up to 40 points", "299,000 UZS / mo", "A hotel with thirty rooms"),
    ],
    "model_side_head": "How they pay",
    "model_side": "Individuals pay by card, through Payme.\n\nVenues and companies pay by invoice and bank transfer — which is how a business here actually pays.",
    "nav_math": "The arithmetic",
    "math_title": "Recurring revenue, at the smallest band.",
    "scenarios": [
        ("100 venues", "14.9M UZS", "≈ 179M UZS a year"),
        ("500 venues", "74.5M UZS", "≈ 894M UZS a year"),
        ("1,000 venues", "149M UZS", "≈ 1.79B UZS a year"),
    ],
    "math_note": "Monthly, at the entry band of 149,000 UZS and nothing above it — a hotel or a larger restaurant pays twice that. Devices and personal numbers are not counted here at all. Assumption stated rather than modelled: these are venue counts, not a market size.",
    "nav_why": "Why now",
    "why": [
        ("The phones are ready",
         "Every phone sold in recent years reads NFC, and after 2020 a QR code on a table needs no explaining to anyone."),
        ("Service is the battleground",
         "Cafes and hotels here no longer compete only on price. How fast somebody answers a table is now a thing a guest notices and posts about."),
        ("The rails exist",
         "Payme and bank transfer make a monthly subscription collectable from a business, which is what turns this from a gadget into revenue."),
    ],
    "why_title": "Three things are true at once.",
    "nav_status": "Where it stands",
    "status_title": "Built, deployed, and running.",
    "status_live": [
        "Personal profiles, devices, QR codes and analytics — live",
        "Venue: menu and service list in three languages — live",
        "A separate QR and NFC tag for every table, with a printable sheet",
        "Guest requests: waiter, bill, housekeeping, reviews",
        "Counter screen with sound, and a screen that stays awake",
        "Monthly report: response time, busiest tables, busiest hours",
        "Subscriptions, invoices and expiry reminders",
        "Telegram sign-in · Payme integrated · Uzbek, Russian, English",
    ],
    "status_next_head": "Open to anyone, today",
    "status_next": [
        "flex.com.uz/NAV001 — a real venue, open it on your phone",
        "Tap the menu, call a waiter, watch it arrive at the till",
        "Payme certification is next, then the first paying venues",
    ],
    "nav_comp": "Competition",
    "comp_title": "Nobody local sells the venue half.",
    "competitors": [
        ("UNQX", "The same three-letter, three-digit format. Personal profiles only. Live in the Fergana Valley.",
         "No menu, no tables, no requests — the recurring half is missing."),
        ("International NFC cards", "Popl and the rest. Personal profiles, sold worldwide.",
         "No local payment, no invoice for an accountant, no Uzbek."),
        ("Paper and pen", "What every cafe uses today. Free, and reprinted every price change.",
         "Cannot tell you which table is waiting, or for how long."),
    ],
    "nav_team": "Team",
    "team_title": "Who is building it.",
    "team": [
        ("Javohir Abrorov", "Founder · AI engineer", "MC LEGAL"),
        ("Zoirjon Abduvohidov", "Web developer", "Tompson School"),
        ("Abdujabbor Ahmadjonov", "Product manager", "Canada"),
    ],
    "nav_ask": "The ask",
    "ask_title": "What we are raising.",
    "ask_amount": "$50,000 – $100,000",
    "ask_stage": "Pre-seed · the product is built and running",
    "use": [
        ("40%", "Devices: cards, rings, bracelets and table tags in stock"),
        ("35%", "Sales and installation for the first 100 venues"),
        ("15%", "Team"),
        ("10%", "Marketing and events"),
    ],
    "contact_head": "Contact",
    "contact": "Javohir Abrorov\nsheyx2772@gmail.com\n+998 97 724 79 99\nflex.com.uz",
}

UZ = {
    "tagline": "Sizning raqamingiz —\nsizning shaxsingiz.",
    "sub": "Noyob raqam va uni ochadigan NFC buyum.\nBir tegish. Ilova o'rnatilmaydi.",
    "event": "ICT WEEK UZBEKISTAN  ·  2026  ·  PRE-SEED",
    "nav_problem": "Muammo",
    "problem_title": "Qog'oz bir marta bosiladi.\nQolgan hammasi o'zgaradi.",
    "problems": [
        ("Vizitka bosilgan kuniyoq\neskiradi",
         "Yangi ish, yangi raqam, yangi kompaniya — quttidagi vizitkalar noto'g'ri bo'lib qoladi. Odamlar baribir uzatadi."),
        ("Menyu — o'tgan oygi\nnarxlarning surati",
         "Kafe qayta bosadi yoki ruchka bilan chizib tashlaydi. Tushda tugagan taom kechgacha hamma stolda turaveradi."),
        ("Stol hech narsa\nso'ray olmaydi",
         "Mehmon qo'l ko'taradi va kimdir qaraguncha kutadi. Kafe esa qaysi stol qancha kutganini hech qachon bilmaydi."),
    ],
    "nav_what": "Flex nima",
    "what_title": "Raqam, buyum va\nuning ortidagi sahifa.",
    "what_body": "Uch harf va uch raqam — MYN042 — manzilga aylanadi. Karta, uzuk, braslet yoki stol ustidagi belgi uni ochadi.",
    "steps": [
        "Mehmon o'z telefonini tegizadi",
        "Sahifa uning o'z tilida ochiladi",
        "Hech kim hech narsa o'rnatmaydi",
    ],
    "nav_products": "Ikki mahsulot",
    "products_title": "Bitta platforma, ikki xil mijoz.",
    "products": [
        ("Shaxsiy", "Umrbod sizniki bo'ladigan raqam", "100 000 so'mdan",
         ["Profil, havolalar, kontaktlar", "Karta, uzuk yoki braslet", "Kamyoblik bo'yicha narx, bir marta"], True),
        ("Obyekt", "Kafe, mehmonxona va do'konlar", "oyiga 149 000 so'mdan",
         ["Menyu va xizmatlar ro'yxati", "Har stolga alohida belgi", "Chaqiruv, kassa ekrani, hisobot"], False),
    ],
    "nav_loop": "Obyekt qanday ishlaydi",
    "loop_title": "Bosilgan menyu qila olmaydigan yarmi.",
    "loop": [
        "Har stolning o'z QR kodi va o'z NFC belgisi bor",
        "Mehmon tegizadi — menyu uning tilida ochiladi",
        "Bir tegish: ofitsiantni chaqiradi, hisob so'raydi yoki izoh qoldiradi",
        "So'rov kassadagi telefonga stol raqami bilan tushadi va ovoz chiqaradi",
        "Hech kim kirmaydi: kassa telefoni bitta havolani ochib, o'shanda turadi",
    ],
    "nav_gets": "Obyekt nima oladi",
    "gets_title": "Dushanba bosiladi, seshanba javob beradi.",
    "gets_body": "Egasi stollarni yozadi, varaqni chop etadi va har stolga bittadan qo'yadi. Har bir kod o'z stol raqamini olib yuradi — so'rov doim qayerdan kelganini aytadi.",
    "nav_model": "Biznes modeli",
    "model_title": "Raqam bir marta sotib olinadi.\nObyekt har oy to'laydi.",
    "bands": [
        ("Shaxsiy raqam", "100 000 so'mdan", "Kombinatsiya kamyobligiga qarab · bir marta"),
        ("Qurilma", "200 000 – 350 000 so'm", "Karta, uzuk yoki braslet · bir marta"),
        ("Obyekt, 15 nuqtagacha", "oyiga 149 000", "O'n ikki stolli kafe"),
        ("Obyekt, 40 nuqtagacha", "oyiga 299 000", "O'ttiz xonali mehmonxona"),
    ],
    "model_side_head": "Qanday to'laydi",
    "model_side": "Jismoniy shaxs — Payme orqali karta bilan.\n\nObyekt va firma — hisob-faktura va bank o'tkazmasi bilan. Bizda biznes aynan shunday to'laydi.",
    "nav_math": "Hisob-kitob",
    "math_title": "Eng past tarifda takrorlanuvchi daromad.",
    "scenarios": [
        ("100 obyekt", "14,9 mln so'm", "≈ yiliga 179 mln"),
        ("500 obyekt", "74,5 mln so'm", "≈ yiliga 894 mln"),
        ("1 000 obyekt", "149 mln so'm", "≈ yiliga 1,79 mlrd"),
    ],
    "math_note": "Oyiga, eng past 149 000 so'mlik tarif bo'yicha va undan yuqorisi hisobga olinmagan — mehmonxona yoki kattaroq restoran ikki barobar to'laydi. Qurilmalar va shaxsiy raqamlar bu yerda umuman sanalmagan. Bu bozor hajmi emas, obyektlar soni bo'yicha hisob.",
    "nav_why": "Nega hozir",
    "why_title": "Uchta narsa bir vaqtda ro'y berdi.",
    "why": [
        ("Telefonlar tayyor",
         "So'nggi yillarda sotilgan har bir telefon NFC o'qiydi, 2020 dan keyin esa stol ustidagi QR kodni hech kimga tushuntirish kerak emas."),
        ("Raqobat — xizmatda",
         "Kafe va mehmonxonalar endi faqat narxda raqobatlashmaydi. Stolga qancha tez javob berilishi — mehmon sezadigan va yozadigan narsa."),
        ("To'lov yo'llari bor",
         "Payme va bank o'tkazmasi oylik obunani bizneslardan yig'ish imkonini beradi — bu buyumni daromadga aylantiradigan narsa."),
    ],
    "nav_status": "Hozirgi holat",
    "status_title": "Qurilgan, joylashtirilgan, ishlayapti.",
    "status_live": [
        "Shaxsiy profil, qurilmalar, QR va statistika — ishlayapti",
        "Obyekt: menyu va xizmatlar ro'yxati, uch tilda — ishlayapti",
        "Har stolga alohida QR va NFC, chop etiladigan varaq bilan",
        "Mehmon so'rovlari: ofitsiant, hisob, tozalash, izohlar",
        "Ovozli kassa ekrani va o'chmaydigan displey",
        "Oylik hisobot: javob vaqti, gavjum stollar va soatlar",
        "Obuna, hisob-faktura va muddat eslatmalari",
        "Telegram orqali kirish · Payme ulangan · o'zbek, rus, ingliz",
    ],
    "status_next_head": "Bugun har kim ochishi mumkin",
    "status_next": [
        "flex.com.uz/NAV001 — haqiqiy obyekt, telefoningizda oching",
        "Menyuni oching, ofitsiant chaqiring, kassaga tushishini ko'ring",
        "Keyingisi — Payme sertifikatsiyasi va birinchi to'lovchi obyektlar",
    ],
    "nav_comp": "Raqobat",
    "comp_title": "Obyekt yarmini mahalliy hech kim sotmaydi.",
    "competitors": [
        ("UNQX", "Xuddi shu uch harf, uch raqam formati. Faqat shaxsiy profil. Farg'ona vodiysida ishlayapti.",
         "Menyu yo'q, stol yo'q, so'rov yo'q — takrorlanuvchi yarmi yetishmaydi."),
        ("Xalqaro NFC kartalar", "Popl va boshqalar. Shaxsiy profil, butun dunyoda.",
         "Mahalliy to'lov yo'q, buxgalterga hisob-faktura yo'q, o'zbek tili yo'q."),
        ("Qog'oz va ruchka", "Bugun har bir kafe shundan foydalanadi. Bepul, va har narx o'zgarganda qayta bosiladi.",
         "Qaysi stol kutayotganini va qancha kutganini ayta olmaydi."),
    ],
    "nav_team": "Jamoa",
    "team_title": "Kim qurayapti.",
    "team": [
        ("Javohir Abrorov", "Asoschi · AI muhandis", "MC LEGAL"),
        ("Zoirjon Abduvohidov", "Veb dasturchi", "Tompson School"),
        ("Abdujabbor Ahmadjonov", "Mahsulot menejeri", "Kanada"),
    ],
    "nav_ask": "So'rov",
    "ask_title": "Nima so'rayapmiz.",
    "ask_amount": "$50 000 – $100 000",
    "ask_stage": "Pre-seed · mahsulot qurilgan va ishlayapti",
    "use": [
        ("40%", "Qurilmalar: karta, uzuk, braslet va stol belgilari zaxirasi"),
        ("35%", "Birinchi 100 obyekt: savdo va o'rnatish"),
        ("15%", "Jamoa"),
        ("10%", "Marketing va tadbirlar"),
    ],
    "contact_head": "Aloqa",
    "contact": "Javohir Abrorov\nsheyx2772@gmail.com\n+998 97 724 79 99\nflex.com.uz",
}


if __name__ == "__main__":
    OUT.mkdir(parents=True, exist_ok=True)
    for lang, dictionary in (("en", EN), ("uz", UZ)):
        build(lang, dictionary)
    print()

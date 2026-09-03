#!/usr/bin/env python3
"""Packs the rendered slides into a .pptx and a .pdf.

Each slide is one full-bleed image, so the deck looks the same on a projector in
a hall as it does here: no font to be missing, no version of PowerPoint to
disagree about a rounded corner, nothing to reflow. The PDF is for sending —
an investor opens it on a phone without owning anything.

Usage: python3 scripts/deck_pack.py uz
"""

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
W, H = Inches(13.333), Inches(7.5)


def pack(lang):
    src = sorted((ROOT / f"deck/render-{lang}").glob("*.png"))
    if not src:
        raise SystemExit(f"No rendered slides for {lang} — run scripts/deck-render.mjs first.")

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    # Downsampled on the way in. The renders are 3840 across so nothing is soft
    # if somebody zooms, but a deck that cannot be emailed is a deck nobody
    # reads — and 1920 is already more than any projector will show.
    tmp = ROOT / f"deck/.pack-{lang}"
    tmp.mkdir(exist_ok=True)

    for path in src:
        small = tmp / f"{path.stem}.jpg"
        Image.open(path).convert("RGB").resize((1920, 1080), Image.LANCZOS).save(
            small, quality=88, optimize=True
        )
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(small), 0, 0, width=W, height=H)

    pptx = ROOT / f"deck/flex-pitch-{lang}.pptx"
    prs.save(pptx)
    for leftover in tmp.iterdir():
        leftover.unlink()
    tmp.rmdir()

    # The same images as a PDF. Downsampled to 1920 across: a 3840-wide slide
    # makes a file too large to email, and nothing shows a deck at that size.
    pages = []
    for path in src:
        img = Image.open(path).convert("RGB")
        pages.append(img.resize((1920, 1080), Image.LANCZOS))

    pdf = ROOT / f"deck/flex-pitch-{lang}.pdf"
    pages[0].save(pdf, save_all=True, append_images=pages[1:], resolution=150.0)

    print(f"  {pptx.name:<24} {len(src)} slides, {pptx.stat().st_size / 1024 / 1024:.1f} MB")
    print(f"  {pdf.name:<24} {len(src)} slides, {pdf.stat().st_size / 1024 / 1024:.1f} MB")


if __name__ == "__main__":
    for lang in (sys.argv[1:] or ["uz", "en"]):
        pack(lang)

#!/usr/bin/env python3
"""Drawing primitives for the PowerPoint deck.

The deck used to be built as HTML, screenshotted, and pasted into PowerPoint one
full-slide image per page. It looked right and could not be touched: no text to
correct, no picture to swap, no box to move. This module is the replacement —
every heading is a text frame, every card is a rounded rectangle, every
photograph is its own picture. Somebody can open the file and change a price.

Everything is positioned in points, 960 x 540, which is the 16:9 slide at
13.333 x 7.5 inches. The numbers therefore read the same as the CSS they came
from.
"""

import copy

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ── the palette, unchanged from the site ──────────────────────────────────
INK = RGBColor(0x0E, 0x0A, 0x1B)
LIME = RGBColor(0xAB, 0xFF, 0x09)
PAPER = RGBColor(0xFA, 0xFA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Ink-ground text: white, then a grey that still passes at body size.
DIM = RGBColor(0xC2, 0xBE, 0xCE)
FAINT = RGBColor(0x8A, 0x84, 0x9B)
CARD_INK = RGBColor(0x18, 0x12, 0x28)
LINE_INK = RGBColor(0x2A, 0x21, 0x3D)

# Paper-ground text.
INK_DIM = RGBColor(0x5A, 0x55, 0x66)
LINE_PAPER = RGBColor(0xE6, 0xE4, 0xEA)
# The lime is 1.2:1 on white and unreadable as text; this is the same hue taken
# down until it passes, and it is the only green used on the light slides.
OLIVE = RGBColor(0x4D, 0x7C, 0x0F)

HEAD = "Space Grotesk"
BODY = "Inter"

W, H = 960.0, 540.0
M = 64.0          # side margin
TOP = 48.0        # first baseline block


# ── low-level helpers ─────────────────────────────────────────────────────
def _alpha(color_elm, pct):
    """Put an alpha channel on a colour element. python-pptx has no API."""
    a = color_elm.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))})
    color_elm.append(a)


def background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def rect(slide, x, y, w, h, fill=None, line=None, radius=None, shape=None):
    kind = shape or (MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE)
    s = slide.shapes.add_shape(kind, Pt(x), Pt(y), Pt(w), Pt(h))
    s.shadow.inherit = False

    if radius:
        # The adjustment is a fraction of the shorter side, not an absolute.
        s.adjustments[0] = min(0.5, radius / min(w, h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)

    s.text_frame.word_wrap = True
    return s


def scrim(slide, x, y, w, h, stops, angle=5400000):
    """A gradient of the ink at varying opacity, for type over a photograph.

    angle is in 60000ths of a degree: 5400000 is top-to-bottom, 0 left-to-right.
    stops is [(position 0-1, alpha percent), ...].
    """
    s = rect(slide, x, y, w, h, fill=INK)
    s.line.fill.background()

    spPr = s.fill._xPr
    for old in spPr.findall(qn("a:solidFill")):
        spPr.remove(old)

    grad = spPr.makeelement(qn("a:gradFill"), {"rotWithShape": "1"})
    lst = grad.makeelement(qn("a:gsLst"), {})
    for pos, alpha in stops:
        gs = lst.makeelement(qn("a:gs"), {"pos": str(int(pos * 100000))})
        clr = gs.makeelement(qn("a:srgbClr"), {"val": "0E0A1B"})
        _alpha(clr, alpha)
        gs.append(clr)
        lst.append(gs)
    grad.append(lst)
    lin = grad.makeelement(qn("a:lin"), {"ang": str(angle), "scaled": "0"})
    grad.append(lin)

    # The gradient has to sit where the solid fill was, before the line.
    ln = spPr.find(qn("a:ln"))
    spPr.insert(list(spPr).index(ln) if ln is not None else len(spPr), grad)
    return s


def picture(slide, path, x, y, w, h, focus=0.5):
    """Insert a picture cropped to fill the box, the way object-fit: cover does.

    focus moves the kept window along the axis that gets cropped: 0.5 is the
    middle, 0 the left or top edge.
    """
    pic = slide.shapes.add_picture(str(path), Pt(x), Pt(y), Pt(w), Pt(h))
    native = pic.image.size[0] / pic.image.size[1]
    box = w / h

    if native > box:                      # wider than the box: trim the sides
        keep = box / native
        cut = 1.0 - keep
        pic.crop_left = cut * focus
        pic.crop_right = cut * (1.0 - focus)
    else:                                 # taller: trim top and bottom
        keep = native / box
        cut = 1.0 - keep
        pic.crop_top = cut * focus
        pic.crop_bottom = cut * (1.0 - focus)

    pic.left, pic.top, pic.width, pic.height = Pt(x), Pt(y), Pt(w), Pt(h)
    return pic


def contain(slide, path, x, y, w, h):
    """Insert a picture whole, centred in the box, nothing cropped."""
    from PIL import Image

    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    return slide.shapes.add_picture(
        str(path), Pt(x + (w - dw) / 2), Pt(y + (h - dh) / 2), Pt(dw), Pt(dh)
    )


def icon(slide, root, name, tone, x, y, size=22.0):
    path = root / "deck" / "icons-png" / f"{name}-{tone}.png"
    return slide.shapes.add_picture(str(path), Pt(x), Pt(y), Pt(size), Pt(size))


def text(slide, body, x, y, w, h=None, size=16, font=BODY, color=DIM,
         bold=False, align=PP_ALIGN.LEFT, spacing=1.35, tracking=None,
         caps=False, anchor=MSO_ANCHOR.TOP, space_after=0):
    """A text frame. `body` may be a string or a list of (text, overrides)."""
    box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h if h else size * 2.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = body if isinstance(body, list) else [body]
    for i, line in enumerate(lines):
        over = {}
        if isinstance(line, tuple):
            line, over = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = over.get("align", align)
        p.line_spacing = over.get("spacing", spacing)
        p.space_after = Pt(over.get("space_after", space_after))

        run = p.add_run()
        run.text = line.upper() if over.get("caps", caps) else line
        f = run.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)

        track = over.get("tracking", tracking)
        if track is not None:
            # Letter spacing lives on the run properties and has no API.
            run.font._rPr.set("spc", str(int(track * 100)))
    return box


def eyebrow(slide, label, x, y, dark=True, rule=True):
    """The small capitalised label with a lime rule, on every content slide."""
    if rule:
        bar = rect(slide, x, y + 5.5, 18, 2.5, fill=LIME)
        bar.line.fill.background()
    text(slide, label, x + (26 if rule else 0), y, 400, 16,
         size=9, font=BODY, bold=True, tracking=1.6, caps=True,
         color=FAINT if dark else INK_DIM)


def heading(slide, title, x, y, w, dark=True, size=30):
    text(slide, title.split("\n"), x, y, w, None,
         size=size, font=HEAD, bold=True, spacing=1.06,
         color=WHITE if dark else INK, tracking=-0.6)


def footer(slide, page, dark=True):
    tone = FAINT if dark else INK_DIM
    text(slide, "flex.com.uz", M, H - 34, 200, 14, size=8.5, color=tone, tracking=0.4)
    text(slide, str(page), W - M - 60, H - 34, 60, 14, size=8.5, color=tone,
         align=PP_ALIGN.RIGHT)


def bullet(slide, body, x, y, w, dark=False, size=10.5, color=None):
    dot = text(slide, "·", x, y - 2, 10, 16, size=14, bold=True,
               color=LIME if dark else OLIVE)
    text(slide, body, x + 11, y, w - 11, None, size=size,
         color=color or (DIM if dark else INK_DIM), spacing=1.4)
    return dot

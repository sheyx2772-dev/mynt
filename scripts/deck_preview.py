#!/usr/bin/env python3
"""Renders the generated .pptx back to PNGs, so the layout can be checked.

Not a PowerPoint clone — it reads the shapes actually in the file and draws
them: fills, lines, pictures with their crops, gradients, and text wrapped at
the real box width. That is enough to catch what goes wrong when a deck is
built blind — a line that overflows its card, two boxes that overlap, a caption
that runs off the bottom of the slide — which is the whole reason it exists.

PowerPoint's own AppleScript export needs automation permission this machine has
not granted, so this stands in for it.

    python3 scripts/deck_preview.py uz
"""

import io
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
SCALE = 2
W, H = 960 * SCALE, 540 * SCALE

# The deck's own faces, so the preview measures the same widths PowerPoint will.
# Helvetica stands in if they are not installed — close enough in width to show
# whether a line fits, and it errs wide, which catches overflow rather than
# hiding it.
HOME = Path.home() / "Library" / "Fonts"
FONTS = {
    ("Space Grotesk", False): HOME / "SpaceGrotesk-Regular.ttf",
    ("Space Grotesk", True): HOME / "SpaceGrotesk-Bold.ttf",
    ("Inter", False): HOME / "Inter-Regular.ttf",
    ("Inter", True): HOME / "Inter-Bold.ttf",
}
FALLBACK = {
    False: "/System/Library/Fonts/Helvetica.ttc",
    True: "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
}
_cache = {}
missing = set()


def font(size, bold, name="Inter"):
    key = (round(size * SCALE), bold, name)
    if key not in _cache:
        path = FONTS.get((name, bold))
        if path is None or not path.exists():
            missing.add(name)
            path = FALLBACK[bold]
        try:
            _cache[key] = ImageFont.truetype(str(path), key[0])
        except OSError:
            _cache[key] = ImageFont.load_default()
    return _cache[key]


def pt(v):
    return float(Emu(v).pt) * SCALE


def solid_rgb(fill_elm):
    clr = fill_elm.find(qn("a:solidFill"))
    if clr is None:
        return None
    srgb = clr.find(qn("a:srgbClr"))
    if srgb is None:
        return None
    v = srgb.get("val")
    return tuple(int(v[i:i + 2], 16) for i in (0, 2, 4))


def gradient_stops(sp_elm):
    grad = sp_elm.find(qn("a:gradFill"))
    if grad is None:
        return None
    stops = []
    for gs in grad.find(qn("a:gsLst")):
        srgb = gs.find(qn("a:srgbClr"))
        alpha = srgb.find(qn("a:alpha"))
        stops.append((
            int(gs.get("pos")) / 100000,
            tuple(int(srgb.get("val")[i:i + 2], 16) for i in (0, 2, 4)),
            (int(alpha.get("val")) / 100000) if alpha is not None else 1.0,
        ))
    lin = grad.find(qn("a:lin"))
    angle = int(lin.get("ang")) if lin is not None else 5400000
    return stops, angle


def draw_gradient(img, box, stops, angle):
    x, y, w, h = [int(v) for v in box]
    if w <= 0 or h <= 0:
        return
    vertical = angle == 5400000
    n = h if vertical else w
    band = Image.new("RGBA", (1, n) if vertical else (n, 1))
    px = band.load()

    for i in range(n):
        pos = i / max(1, n - 1)
        lo = stops[0]
        hi = stops[-1]
        for a, b in zip(stops, stops[1:]):
            if a[0] <= pos <= b[0]:
                lo, hi = a, b
                break
        span = (hi[0] - lo[0]) or 1
        k = (pos - lo[0]) / span
        colour = tuple(int(lo[1][c] + (hi[1][c] - lo[1][c]) * k) for c in range(3))
        alpha = int((lo[2] + (hi[2] - lo[2]) * k) * 255)
        if vertical:
            px[0, i] = colour + (alpha,)
        else:
            px[i, 0] = colour + (alpha,)

    img.alpha_composite(band.resize((w, h)), (x, y))


def wrap(draw, body, f, width):
    lines = []
    for para in body.split("\n"):
        words, line = para.split(" "), ""
        for word in words:
            trial = f"{line} {word}".strip()
            if draw.textlength(trial, font=f) <= width or not line:
                line = trial
            else:
                lines.append(line)
                line = word
        lines.append(line)
    return lines


def render(lang):
    src = ROOT / "deck" / f"flex-pitch-{lang}.pptx"
    out = ROOT / "deck" / f"preview-{lang}"
    out.mkdir(parents=True, exist_ok=True)
    for old in out.glob("*.png"):
        old.unlink()

    prs = Presentation(str(src))
    for index, slide in enumerate(prs.slides, start=1):
        ground = solid_rgb(slide.background.fill._xPr) or (255, 255, 255)
        img = Image.new("RGBA", (W, H), ground + (255,))
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            box = (pt(shape.left), pt(shape.top), pt(shape.width), pt(shape.height))
            x, y, w, h = box

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
                iw, ih = blob.size
                left = int(iw * shape.crop_left)
                top = int(ih * shape.crop_top)
                right = iw - int(iw * shape.crop_right)
                bottom = ih - int(ih * shape.crop_bottom)
                if right > left and bottom > top:
                    blob = blob.crop((left, top, right, bottom))
                if w >= 1 and h >= 1:
                    img.alpha_composite(
                        blob.resize((int(w), int(h)), Image.LANCZOS), (int(x), int(y))
                    )
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                sp = shape.fill._xPr
                grad = gradient_stops(sp)
                if grad:
                    draw_gradient(img, box, *grad)
                else:
                    rgb = solid_rgb(sp)
                    if rgb:
                        draw.rounded_rectangle(
                            [x, y, x + w, y + h],
                            radius=min(14 * SCALE, w / 2, h / 2),
                            fill=rgb + (255,),
                        )
                    ln = shape._element.spPr.find(qn("a:ln"))
                    line = solid_rgb(ln) if ln is not None else None
                    if line:
                        draw.rounded_rectangle(
                            [x, y, x + w, y + h],
                            radius=min(14 * SCALE, w / 2, h / 2),
                            outline=line + (255,), width=max(1, SCALE),
                        )

            if not shape.has_text_frame:
                continue

            ty = y
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    f = font(run.font.size.pt if run.font.size else 12,
                             bool(run.font.bold), run.font.name or "Inter")
                    colour = (0, 0, 0)
                    try:
                        colour = tuple(run.font.color.rgb)
                    except Exception:
                        pass
                    spacing = para.line_spacing or 1.3
                    step = (run.font.size.pt if run.font.size else 12) * spacing * SCALE
                    for line in wrap(draw, run.text, f, max(10, w)):
                        tx = x
                        if str(para.alignment) == "CENTER (2)":
                            tx = x + (w - draw.textlength(line, font=f)) / 2
                        elif str(para.alignment) == "RIGHT (3)":
                            tx = x + w - draw.textlength(line, font=f)
                        draw.text((tx, ty), line, font=f, fill=colour + (255,))
                        ty += step
                    ty += (para.space_after.pt * SCALE) if para.space_after else 0

        img.convert("RGB").save(out / f"{index:02d}.png", quality=92)

    note = f"  (substituted: {', '.join(sorted(missing))})" if missing else ""
    print(f"  {len(prs.slides._sldIdLst)} slides → deck/preview-{lang}/{note}")


if __name__ == "__main__":
    for lang in (sys.argv[1:] or ["uz"]):
        render(lang)

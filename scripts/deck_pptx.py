#!/usr/bin/env python3
"""Builds the pitch deck as a real, editable PowerPoint file.

Nineteen slides, and on every one of them the heading is a text frame, the cards
are rounded rectangles and the photographs are pictures. Open it and change a
price, swap a photograph, move a box. The earlier version put one flat image on
each slide: it looked the same and could not be edited at all, which for a deck
that goes to an event and then round a dozen inboxes is the wrong trade.

The fonts are Space Grotesk and Inter, both free from Google Fonts. If they are
not installed PowerPoint substitutes something and the layout still holds; with
them installed the file is the design as drawn.

    python3 scripts/deck_pptx.py uz en
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import deck_notes
from deck_lib import (
    BODY, CARD_INK, DIM, FAINT, HEAD, INK, INK_DIM, LIME, LINE_INK, LINE_PAPER,
    M, OLIVE, PAPER, TOP, W, WHITE, H,
    background, bullet, contain, eyebrow, footer, heading, icon, picture,
    rect, scrim, text,
)

ROOT = Path(__file__).resolve().parent.parent
PHOTOS = ROOT / "deck" / "graded"


def photo(name):
    return PHOTOS / f"{name}.jpg"


def card(slide, x, y, w, h, dark):
    return rect(slide, x, y, w, h, radius=14,
                fill=CARD_INK if dark else WHITE,
                line=LINE_INK if dark else LINE_PAPER)


def page(prs):
    """The slide's own number. Written by hand once, and wrong by the evening."""
    return len(prs.slides._sldIdLst) - 1


def blank(prs, dark=True, ground=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide, ground if ground is not None else (INK if dark else PAPER))
    return slide


# ── 1 · cover ──────────────────────────────────────────────────────────────
def cover(prs, t):
    s = blank(prs)
    picture(s, photo("qurilmalar"), 0, 0, W, H, focus=0.5)
    scrim(s, 0, 0, W, H, [(0.0, 97), (0.34, 86), (0.52, 30), (0.68, 0)])

    eyebrow(s, t["event"], M, TOP)
    text(s, "flex", M, TOP + 26, 400, 74, size=56, font=HEAD, bold=True,
         color=WHITE, tracking=-2.0)
    text(s, t["tagline"], M, TOP + 100, 560, None, size=27, font=HEAD,
         bold=True, color=LIME, spacing=1.1, tracking=-0.6)
    text(s, t["sub"], M, TOP + 166, 470, None, size=11.5, color=DIM, spacing=1.5)


# ── 2 · problem ────────────────────────────────────────────────────────────
def problem(prs, t):
    s = blank(prs)
    eyebrow(s, t["nav_problem"], M, TOP)
    heading(s, t["problem_title"], M, TOP + 22, 700)

    cw, ch, gap = 200.0, 132.0, 10.5
    x0, y0 = M, 190.0
    for i, (ic, head, bodycopy) in enumerate(t["problems"]):
        x = x0 + (i % 4) * (cw + gap)
        y = y0 + (i // 4) * (ch + gap)
        card(s, x, y, cw, ch, dark=True)
        icon(s, ROOT, ic, "lime", x + 18, y + 16, 20)
        text(s, head, x + 18, y + 44, cw - 34, None, size=11.5, font=HEAD,
             bold=True, color=WHITE, spacing=1.15)
        text(s, bodycopy, x + 18, y + 70, cw - 34, None, size=8.6,
             color=FAINT, spacing=1.42)
    footer(s, page(prs))


# ── 3 · the thesis, alone, in lime ─────────────────────────────────────────
def statement(prs, t):
    s = blank(prs, ground=LIME)
    text(s, t["statement"].replace("<br>", "\n").split("\n"),
         M, 176, 800, None, size=48, font=HEAD, bold=True, color=INK,
         spacing=1.04, tracking=-1.6)
    text(s, t["statement_sub"], M, 320, 540, None, size=13,
         color=INK_DIM, spacing=1.5)


# ── 4 · every problem, answered ────────────────────────────────────────────
def answer(prs, t):
    s = blank(prs, dark=False)
    eyebrow(s, t["nav_answer"], M, TOP, dark=False)
    heading(s, t["answer_title"], M, TOP + 22, 430, dark=False)
    text(s, t["answer_body"], 520, TOP + 30, 376, None, size=10.5,
         color=INK_DIM, spacing=1.5)

    cw, rh = 400.0, 44.0
    x0, y0 = M, 178.0
    for i, (problem_text, fix) in enumerate(t["answers"]):
        x = x0 + (i % 2) * (cw + 32)
        y = y0 + (i // 2) * (rh + 8)
        rect(s, x, y + rh - 1, cw, 0.75, fill=LINE_PAPER)
        text(s, problem_text, x, y + 4, 150, None, size=10, font=HEAD,
             bold=True, color=INK, spacing=1.25)
        text(s, fix, x + 160, y + 4, cw - 160, None, size=9.6,
             color=INK_DIM, spacing=1.35)
    footer(s, page(prs), dark=False)


# ── 5 · the hardware ───────────────────────────────────────────────────────
TRIMMED = {"karta", "uzuk", "braslet"}

# Where in the frame the product actually is. The tag on the dog's collar sits
# low in its picture, and a centred crop showed a portrait of a dog.
FOCUS = {"hayvon-teg": 0.86, "buyum-teg": 0.55, "avtovizitka": 0.5}


def devices(prs, t):
    s = blank(prs, dark=False)
    eyebrow(s, t["nav_devices"], M, TOP, dark=False)
    heading(s, t["devices_title"], M, TOP + 22, 620, dark=False)

    cw, ch, gap = 199.0, 157.0, 12.0
    x0, y0 = M, 172.0
    for i, (key, name, where, price) in enumerate(t["devices"]):
        x = x0 + (i % 4) * (cw + gap)
        y = y0 + (i // 4) * (ch + gap)
        card(s, x, y, cw, ch, dark=False)

        band = 86.0
        if key in TRIMMED:
            contain(s, photo(key), x + 8, y + 6, cw - 16, band - 6)
        else:
            picture(s, photo(key), x + 8, y + 8, cw - 16, band - 14,
                    focus=FOCUS.get(key, 0.5))

        text(s, name, x + 14, y + band + 9, cw - 26, None, size=11,
             font=HEAD, bold=True, color=INK)
        text(s, where, x + 14, y + band + 26, cw - 26, None, size=8.0,
             color=INK_DIM, spacing=1.35)
        text(s, price, x + 14, y + ch - 24, cw - 26, None, size=9.6,
             font=HEAD, bold=True, color=OLIVE)
    footer(s, page(prs), dark=False)


# ── 6 · why it is not an ordinary NFC card ────────────────────────────────
def notplain(prs, t):
    s = blank(prs)
    eyebrow(s, t["nav_notplain"], M, TOP)
    heading(s, t["notplain_title"], M, TOP + 22, 430)
    text(s, t["notplain_body"], 520, TOP + 30, 376, None, size=10.5,
         color=FAINT, spacing=1.5)

    cw, ch, gap = 200.0, 126.0, 10.5
    x0, y0 = M, 156.0
    for i, (ic, head, bodycopy) in enumerate(t["notplain"]):
        x = x0 + (i % 4) * (cw + gap)
        y = y0 + (i // 4) * (ch + gap)
        card(s, x, y, cw, ch, dark=True)
        icon(s, ROOT, ic, "lime", x + 18, y + 15, 19)
        text(s, head, x + 18, y + 41, cw - 34, None, size=11, font=HEAD,
             bold=True, color=WHITE, spacing=1.15)
        text(s, bodycopy, x + 18, y + 66, cw - 34, None, size=8.4,
             color=FAINT, spacing=1.42)
    footer(s, page(prs))


# ── 7 · card designs ───────────────────────────────────────────────────────
def designs(prs, t):
    s = blank(prs, dark=False)
    eyebrow(s, t["nav_design"], M, TOP, dark=False)
    heading(s, t["design_title"], M, TOP + 22, 620, dark=False)
    text(s, t["design_body"], M, 138, 640, None, size=11.5,
         color=INK_DIM, spacing=1.55)

    cw, ch, gap = 200.0, 128.0, 10.5
    y = 246.0
    for i, (key, name) in enumerate(t["card_designs"]):
        x = M + i * (cw + gap)
        picture(s, photo(key), x, y, cw, ch)
        text(s, name, x, y + ch + 12, cw, 18, size=10.5, font=HEAD,
             bold=True, color=INK)
    footer(s, page(prs), dark=False)


# ── 8 · the market, in photographs ────────────────────────────────────────
def market(prs, t):
    s = blank(prs)
    eyebrow(s, t["nav_market"], M, TOP)
    heading(s, t["market_title"], M, TOP + 22, 700)

    top = 148.0
    cw = (W - 3 * 2) / 4
    for i, (key, ic, name, what) in enumerate(t["verticals"]):
        x = i * (cw + 2)
        picture(s, photo(key), x, top, cw, H - top, focus=0.5)
        scrim(s, x, top + (H - top) * 0.36, cw, (H - top) * 0.64,
              [(0.0, 0), (0.42, 74), (1.0, 95)])
        icon(s, ROOT, ic, "lime", x + 22, H - 118, 20)
        text(s, name, x + 22, H - 90, cw - 40, None, size=12.5, font=HEAD,
             bold=True, color=WHITE)
        text(s, what, x + 22, H - 70, cw - 40, None, size=8.6,
             color=DIM, spacing=1.4)


# ── 9 · market size, in soum ──────────────────────────────────────────────
def size(prs, t):
    s = blank(prs, dark=False)
    eyebrow(s, t["nav_size"], M, TOP, dark=False)
    heading(s, t["size_title"], M, TOP + 22, 620, dark=False)

    rw, rh = 560.0, 76.0
    y0 = 148.0
    for i, (tag, count, money, note) in enumerate(t["size_rows"]):
        y = y0 + i * (rh + 10)
        card(s, M, y, rw, rh, dark=False)
        rect(s, M, y, 4.5, rh, fill=LIME if i == 0 else LINE_PAPER)
        text(s, tag, M + 20, y + 15, 60, 20, size=10, font=HEAD, bold=True,
             color=OLIVE, tracking=1.0)
        text(s, money, M + 20, y + 32, 260, None, size=19, font=HEAD,
             bold=True, color=INK, tracking=-0.5)
        text(s, count, M + 300, y + 15, 240, 18, size=10, font=HEAD,
             bold=True, color=INK)
        text(s, note, M + 300, y + 34, 244, None, size=8.4, color=INK_DIM,
             spacing=1.35)

    x = M + rw + 34
    for i, (ic, number, label) in enumerate(t["size_units"]):
        y = y0 + i * 60
        icon(s, ROOT, ic, "olive", x, y + 4, 17)
        text(s, number, x + 26, y, 220, 26, size=16, font=HEAD, bold=True,
             color=INK, tracking=-0.4)
        text(s, label, x + 26, y + 22, 224, None, size=8.2, color=INK_DIM,
             spacing=1.35)

    text(s, t["size_source"], M, H - 68, 700, None, size=7.8, color=INK_DIM,
         spacing=1.4)
    footer(s, page(prs), dark=False)


# ── 10 · competition ──────────────────────────────────────────────────────
def rivals(prs, t):
    s = blank(prs)
    eyebrow(s, t["nav_rivals"], M, TOP)
    heading(s, t["rivals_title"], M, TOP + 22, 620)

    name_w = 168.0
    col_w = (W - 2 * M - name_w) / 5
    y0 = 150.0

    for j, col in enumerate(t["rivals_cols"]):
        text(s, col, M + name_w + j * col_w, y0, col_w - 10, 26, size=8,
             font=BODY, bold=True, caps=True, tracking=1.0, color=FAINT,
             spacing=1.25)

    y = y0 + 34
    for name, cells, mine in t["rivals"]:
        rh = 46.0
        if mine:
            rect(s, M - 12, y - 6, W - 2 * M + 24, rh, radius=10,
                 fill=CARD_INK, line=LINE_INK)
        text(s, name, M, y + 6, name_w - 14, None, size=11,
             font=HEAD, bold=True, color=LIME if mine else WHITE)
        for j, cell in enumerate(cells):
            text(s, cell, M + name_w + j * col_w, y + 7, col_w - 12, None,
                 size=8.4, color=DIM if mine else FAINT, spacing=1.35)
        if not mine:
            rect(s, M, y + rh - 6, W - 2 * M, 0.75, fill=LINE_INK)
        y += rh

    text(s, t["rivals_note"], M, H - 74, 760, None, size=9, color=FAINT,
         spacing=1.45)
    footer(s, page(prs))


# ── 11 · how the venue half works ─────────────────────────────────────────
def loop(prs, t):
    s = blank(prs)
    eyebrow(s, t["nav_loop"], M, TOP)
    heading(s, t["loop_title"], M, TOP + 22, 620)

    y = 150.0
    picture(s, photo("restoran"), M, y, 186, 300, focus=0.5)
    frame = ROOT / "deck" / "frames" / "counter.png"
    contain(s, frame, M + 202, y - 6, 178, 312)

    x = M + 400
    for i, step in enumerate(t["loop"]):
        ty = y + i * 56
        n = rect(s, x, ty, 24, 24, radius=8, fill=LIME)
        text(s, str(i + 1), x, ty + 5, 24, 16, size=10, font=HEAD, bold=True,
             color=INK, align=PP_ALIGN.CENTER)
        text(s, step, x + 36, ty + 3, W - M - x - 36, None, size=10,
             color=WHITE, spacing=1.4)
    footer(s, page(prs))


# ── 12 · what a venue receives ────────────────────────────────────────────
def gets(prs, t):
    s = blank(prs, dark=False)
    eyebrow(s, t["nav_gets"], M, TOP, dark=False)
    heading(s, t["gets_title"], M, TOP + 22, 700, dark=False)

    y = 158.0
    picture(s, ROOT / "deck" / "shots" / "cards.png", M, y, 420, 244, focus=0.0)
    text(s, t["gets_body"], M, y + 258, 420, None, size=9.4, color=INK_DIM,
         spacing=1.45)

    x = M + 452
    for i, (ic, label) in enumerate(t["gets_items"]):
        ty = y + i * 50
        card(s, x, ty, W - M - x, 40, dark=False)
        icon(s, ROOT, ic, "olive", x + 16, ty + 11, 18)
        text(s, label, x + 44, ty + 13, W - M - x - 58, None, size=9.6,
             font=HEAD, bold=True, color=INK)
    footer(s, page(prs), dark=False)


# ── 13 · business model ───────────────────────────────────────────────────
def model(prs, t):
    s = blank(prs)
    eyebrow(s, t["nav_model"], M, TOP)
    heading(s, t["model_title"], M, TOP + 22, 620)

    y0 = 176.0
    for i, (label, price, note) in enumerate(t["bands"]):
        y = y0 + i * 62
        text(s, label, M, y, 300, None, size=11.5, font=HEAD, bold=True,
             color=WHITE)
        text(s, note, M, y + 20, 300, None, size=8.4, color=FAINT,
             spacing=1.35)
        text(s, price, M + 320, y + 2, 240, 30, size=15, font=HEAD, bold=True,
             color=LIME, align=PP_ALIGN.RIGHT, tracking=-0.3)
        rect(s, M, y + 46, 560, 0.75, fill=LINE_INK)

    x = M + 606
    card(s, x, y0 - 8, W - M - x, 254, dark=True)
    icon(s, ROOT, "file-text", "lime", x + 22, y0 + 16, 22)
    text(s, t["model_side_head"], x + 22, y0 + 48, W - M - x - 44, None,
         size=12, font=HEAD, bold=True, color=WHITE)
    text(s, t["model_side"], x + 22, y0 + 74, W - M - x - 44, None, size=9.2,
         color=FAINT, spacing=1.5)
    footer(s, page(prs))


# ── 14 · the arithmetic ───────────────────────────────────────────────────
def arithmetic(prs, t):
    s = blank(prs, dark=False)
    eyebrow(s, t["nav_math"], M, TOP, dark=False)
    heading(s, t["math_title"], M, TOP + 22, 700, dark=False)

    cw, gap = 266.0, 14.0
    y = 168.0
    for i, (label, number, note) in enumerate(t["scenarios"]):
        x = M + i * (cw + gap)
        card(s, x, y, cw, 210, dark=False)
        text(s, label, x + 24, y + 26, cw - 48, 18, size=8.2, bold=True,
             caps=True, tracking=1.2, color=INK_DIM)
        text(s, number, x + 24, y + 62, cw - 48, 58, size=42, font=HEAD,
             bold=True, color=INK, tracking=-1.6)
        text(s, note, x + 24, y + 132, cw - 48, None, size=9.5, color=INK_DIM,
             spacing=1.4)
        rect(s, x + 24, y + 174, 34, 2.5, fill=LIME)

    text(s, t["math_note"], M, y + 236, 780, None, size=9.2, color=INK_DIM,
         spacing=1.5)
    footer(s, page(prs), dark=False)


# ── 15 · who we are ───────────────────────────────────────────────────────
def traction(prs, t):
    s = blank(prs)
    eyebrow(s, t["nav_traction"], M, TOP)
    heading(s, t["traction_title"], M, TOP + 22, 620)

    cw, ch, gap = 266.0, 132.0, 14.0
    x0, y0 = M, 178.0
    for i, (ic, head, bodycopy) in enumerate(t["traction"]):
        x = x0 + (i % 3) * (cw + gap)
        y = y0 + (i // 3) * (ch + gap)
        card(s, x, y, cw, ch, dark=True)
        icon(s, ROOT, ic, "lime", x + 20, y + 18, 22)
        text(s, head, x + 20, y + 50, cw - 40, None, size=12, font=HEAD,
             bold=True, color=WHITE)
        text(s, bodycopy, x + 20, y + 74, cw - 40, None, size=9,
             color=FAINT, spacing=1.45)
    footer(s, page(prs))


# ── 16 · where it stands, with the one real photograph ───────────────────
def status(prs, t):
    s = blank(prs, dark=False)
    eyebrow(s, t["nav_status"], M, TOP, dark=False)
    heading(s, t["status_title"], M, TOP + 22, 700, dark=False)

    y = 150.0
    for i, line in enumerate(t["status_live"]):
        bullet(s, line, M, y + i * 30, 430, size=9.6)

    x = M + 470
    card(s, x, y - 4, 190, 300, dark=False)
    rect(s, x, y - 4, 190, 300, radius=14, fill=LIME)
    icon(s, ROOT, "smartphone", "ink", x + 22, y + 22, 22)
    text(s, t["status_next_head"], x + 22, y + 58, 148, None, size=11.5,
         font=HEAD, bold=True, color=INK, spacing=1.2)
    text(s, t["status_next"], x + 22, y + 98, 148, None, size=8.8,
         color=INK_DIM, spacing=1.5)

    px = x + 214
    picture(s, photo("haqiqiy-kafe"), px, y - 4, W - M - px, 278, focus=0.62)
    text(s, t["status_photo_caption"], px, y + 284, W - M - px, 16,
         size=7.6, bold=True, caps=True, tracking=1.1, color=INK_DIM)
    footer(s, page(prs), dark=False)


# ── 17 · roadmap ──────────────────────────────────────────────────────────
def roadmap(prs, t):
    s = blank(prs, dark=False)
    eyebrow(s, t["nav_road"], M, TOP, dark=False)
    heading(s, t["road_title"], M, TOP + 22, 620, dark=False)

    cw = (W - 2 * M - 3 * 14) / 4
    y = 168.0
    rect(s, M, y + 13, W - 2 * M, 1.5, fill=LINE_PAPER)

    for i, (when, head, items) in enumerate(t["road"]):
        x = M + i * (cw + 14)
        dot = rect(s, x, y + 8, 12, 12, radius=6,
                   fill=LIME if i == 0 else WHITE, line=LINE_PAPER)
        text(s, when, x, y + 34, cw - 8, 18, size=8.4, bold=True, caps=True,
             tracking=1.1, color=OLIVE)
        text(s, head, x, y + 54, cw - 8, None, size=13, font=HEAD, bold=True,
             color=INK, spacing=1.15)
        for j, item in enumerate(items):
            bullet(s, item, x, y + 88 + j * 42, cw - 8, size=9)

    text(s, t["road_note"], M, H - 66, 600, None, size=8.4, color=INK_DIM)
    footer(s, page(prs), dark=False)


# ── 17 · why now ──────────────────────────────────────────────────────────
def why(prs, t):
    s = blank(prs)
    eyebrow(s, t["nav_why"], M, TOP)
    heading(s, t["why_title"], M, TOP + 22, 620)

    cw, gap = 266.0, 14.0
    y = 186.0
    for i, (ic, head, bodycopy) in enumerate(t["why"]):
        x = M + i * (cw + gap)
        card(s, x, y, cw, 190, dark=True)
        icon(s, ROOT, ic, "lime", x + 20, y + 20, 22)
        text(s, head, x + 20, y + 54, cw - 40, None, size=12.5, font=HEAD,
             bold=True, color=WHITE)
        text(s, bodycopy, x + 20, y + 80, cw - 40, None, size=9,
             color=FAINT, spacing=1.45)
    footer(s, page(prs))


# ── 18 · the team ─────────────────────────────────────────────────────────
def team(prs, t):
    s = blank(prs, dark=False)
    eyebrow(s, t["nav_team"], M, 176, dark=False)
    heading(s, t["team_title"], M, 198, 330, dark=False)
    text(s, t["team_body"], M, 262, 330, None, size=11, color=INK_DIM,
         spacing=1.55)

    x, y0 = 460.0, 158.0
    for i, (initials, name, role, where) in enumerate(t["team"]):
        y = y0 + i * 78
        card(s, x, y, W - M - x, 64, dark=False)
        chip = rect(s, x + 18, y + 14, 36, 36, radius=11, fill=LIME)
        text(s, initials, x + 18, y + 24, 36, 18, size=11, font=HEAD,
             bold=True, color=INK, align=PP_ALIGN.CENTER)
        text(s, name, x + 66, y + 17, 240, None, size=11.5, font=HEAD,
             bold=True, color=INK)
        text(s, role, x + 66, y + 36, 240, None, size=9, bold=True,
             color=OLIVE)
        text(s, where, W - M - 150, y + 26, 132, None, size=9,
             color=INK_DIM, align=PP_ALIGN.RIGHT)
    footer(s, page(prs), dark=False)


# ── 19 · the ask ──────────────────────────────────────────────────────────
def ask(prs, t):
    s = blank(prs)
    picture(s, photo("oila"), 0, 0, W, H, focus=0.6)
    scrim(s, 0, 0, W, H, [(0.0, 92), (0.55, 82), (1.0, 70)], angle=0)

    eyebrow(s, t["nav_ask"], M, TOP)
    text(s, t["ask_amount"], M, TOP + 26, 620, 62, size=44, font=HEAD,
         bold=True, color=LIME, tracking=-1.4)
    text(s, t["ask_stage"], M, TOP + 88, 520, None, size=11, color=DIM)

    y = 190.0
    for i, (share, what) in enumerate(t["use"]):
        ty = y + i * 46
        text(s, share, M, ty, 60, 22, size=13, font=HEAD, bold=True,
             color=LIME)
        text(s, what, M + 70, ty + 2, 430, None, size=10, color=WHITE,
             spacing=1.4)

    x = W - M - 250
    card(s, x, 178, 250, 190, dark=True)
    text(s, t["contact_head"], x + 22, 200, 210, 22, size=12, font=HEAD,
         bold=True, color=WHITE)
    text(s, t["contact"].replace("<br>", "\n").split("\n"),
         x + 22, 230, 210, None, size=9.6, color=DIM, spacing=1.7)


SLIDES = [cover, problem, statement, answer, devices, notplain, designs,
          market, size, rivals, loop, gets, model, arithmetic, traction,
          status, roadmap, why, team, ask]


def build(lang):
    t = json.loads((ROOT / "deck" / f"content-{lang}.json").read_text())

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    for make in SLIDES:
        make(prs, t)

    # What to say over each slide, in the file's own speaker notes. A deck read
    # aloud from its bullet points is the commonest way a good product pitches
    # badly; the words belong here, where only the presenter sees them.
    notes = deck_notes.UZ if lang == "uz" else deck_notes.EN
    for index, (slide, line) in enumerate(zip(prs.slides, notes), start=1):
        mark = "  ·  [3 DAQIQA]" if index in deck_notes.SHORT else ""
        if lang != "uz":
            mark = mark.replace("3 DAQIQA", "3-MINUTE CUT")
        slide.notes_slide.notes_text_frame.text = f"{index}/{len(notes)}{mark}\n\n{line}"

    out = ROOT / "deck" / f"flex-pitch-{lang}.pptx"
    prs.save(out)
    print(f"  {out.name:<24} {len(SLIDES)} slides, "
          f"{out.stat().st_size / 1024 / 1024:.1f} MB, editable")
    return out


if __name__ == "__main__":
    for lang in (sys.argv[1:] or ["uz", "en"]):
        build(lang)
